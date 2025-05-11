# Prima di eseguire questo script, assicurati di installare le librerie necessarie:
# pip install google-generativeai pypdf docx2txt openpyxl unstructured python-pptx

import os
import asyncio
import google.generativeai as genai
import PyPDF2
import docx2txt
from unstructured.partition.xlsx import partition_xlsx
from pptx import Presentation

# Configurazione Gemini
GOOGLE_API_KEY = "AIzaSyASlsh1jRxd2iZnM3E3FM1TlDZ4yarUUMU"
genai.configure(api_key=GOOGLE_API_KEY)

# Configurazione del modello
model = genai.GenerativeModel('gemini-2.0-flash')

async def extract_text_from_docx(file_path):
    """Estrae il testo da un file DOCX."""
    try:
        text = docx2txt.process(file_path)
        return text
    except Exception as e:
        print(f"Errore nell'estrazione del testo da DOCX: {e}")
        return None

async def extract_text_from_pdf(file_path):
    """Estrae il testo da un file PDF."""
    try:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    except Exception as e:
        print(f"Errore nell'estrazione del testo da PDF: {e}")
        return None

async def extract_text_from_xlsx(file_path):
    """Estrae il testo da un file XLSX."""
    try:
        elements = partition_xlsx(file_path)
        text = "\n".join([str(element) for element in elements])
        return text
    except Exception as e:
        print(f"Errore nell'estrazione del testo da XLSX: {e}")
        return None

async def extract_text_from_pptx(file_path):
    """Estrae il testo da un file PowerPoint."""
    try:
        prs = Presentation(file_path)
        text = []
        
        # Estrai testo da ogni slide
        for slide in prs.slides:
            slide_text = []
            # Estrai testo dai placeholder
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            # Aggiungi il testo della slide alla lista
            if slide_text:
                text.append("\n".join(slide_text))
        
        return "\n\n".join(text)
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return None

async def get_document_text(file_path):
    """Estrae il testo dal documento in base al suo tipo."""
    _, file_extension = os.path.splitext(file_path)
    
    if file_extension.lower() == '.docx':
        return await extract_text_from_docx(file_path)
    elif file_extension.lower() == '.pdf':
        return await extract_text_from_pdf(file_path)
    elif file_extension.lower() in ['.xlsx', '.xls']:
        return await extract_text_from_xlsx(file_path)
    elif file_extension.lower() == '.pptx':
        return await extract_text_from_pptx(file_path)
    else:
        print(f"Unsupported file format: {file_extension}")
        return None

async def summarize_text(text, file_type):
    """Genera un riassunto del testo usando Gemini con un prompt specifico per il tipo di file."""
    try:
        # Prompt specifici per tipo di file
        prompts = {
            '.docx': """You are an Expert Document Analyst. Your task is to create a **concise yet faithful textual re-elaboration** of the following text, which is an exact Python transcription from a Word document (`.docx`). Your primary goal is to significantly reduce the length of the input transcription by selectively omitting non-essential details, while meticulously preserving all core messages, key arguments, structural flow of main points, and critical information. The output must be a lean but accurate representation of the original.

            This concise re-elaboration is critically important: it will serve as the primary reference point (representing an older version of the document) for a subsequent AI-driven comparison against a newer version to detect significant content changes.

            Given that the text originates from a `.docx` file, the transcription likely represents narrative content, reports, or structured text. Your concise re-elaboration should reflect this inherent nature.

            Please structure your output as follows:

            1.  **Main Concise Re-elaboration (Flowing, Selective, True to Core Content):**
                * Begin by identifying and briefly stating the document's primary purpose and main scope if clearly discernible from the transcription.
                * Focus on extracting and presenting the **key sections, core arguments, principal findings, and essential conclusions** from the transcription.
                * **To achieve conciseness, you must actively cut redundant phrases, verbose explanations, less critical examples, or overly detailed supporting information.** However, ensure that the remaining text accurately represents the original's main points.
                * Where possible, use key phrasing from the original transcription for the essential elements to maintain faithfulness, but rephrase for brevity when necessary.
                * Preserve the logical flow and apparent structure of the *key information* from the original document. The re-elaboration should still feel like the original document in its core structure, but significantly shorter.
                * The final text should be a continuous, flowing narrative of the essential content, substantially reduced in overall length compared to the input transcription.

            2.  **Keywords (Formatted List - Crucial for Comparison):**
                * The re-elaboration *must* conclude with a list of the most significant and representative keywords and key terms/phrases directly identifiable from the **core content** of the transcribed document text.
                * This list is vital for the subsequent comparison process. Choose keywords that truly encapsulate the essential topics and unique identifiers of the document's core message.
                * Format this list *exactly* as follows, on a new line after the main re-elaboration:
                    `KEYWORDS: keyword1, keyword2, keyword3, term phrase 4, ...`

            ***IMPORTANT: DO NOT include NEVER AND NEVER this character in your response: `, ", ' ***

            DOCUMENT TRANSCRIPTION (from .docx):
            {text}

            CONCISE RE-ELABORATION:""",

            '.pdf': """You are an Expert Document Analyst. Your task is to create a faithful and structured textual re-elaboration of the following text, which is an exact Python transcription from a PDF document (`.pdf`). The goal is to make this re-elaboration as similar as possible to the original transcribed content in terms of information, phrasing, and structure, while ensuring a slightly enhanced text flow and clarity for subsequent AI analysis. **Avoid significant summarization, content alteration, or reinterpretation.**

            This re-elaboration is critically important: it will serve as the primary reference point (representing an older version of the document) for a subsequent AI-driven comparison against a newer version to detect significant content changes.

            Given that the text originates from a `.pdf` file, the transcription might represent a wide variety of content types (reports, articles, forms, etc.). Your re-elaboration should mirror the nature of the provided transcribed text.

            Please structure your output as follows:

            1.  **Main Textual Re-elaboration (Flowing, but True to Original Transcription):**
                * If a clear main purpose or scope is immediately evident from the very beginning of the transcription, you may briefly state it as an introduction. Otherwise, proceed directly to re-elaborating the transcribed content.
                * Present the content, sections (like chapters or distinct textual blocks), and themes largely as they appear in the provided transcription. Focus on preserving the original information, specific phrasing, and apparent structure of the transcription to the maximum extent possible.
                * Ensure that important arguments, findings, specific examples, data points, or any textual descriptions of visual elements (if present in the transcription) are clearly presented, maintaining their original context and detail as found in the input text.
                * You may *lightly* improve readability and logical flow where the transcription might be disjointed or raw. However, **prioritize fidelity to the transcribed content, its order, and its level of detail over extensive rephrasing or narrative creation.** The output should feel very much like the original transcription, just slightly polished for coherence.

            2.  **Keywords (Formatted List - Crucial for Comparison):**
                * The re-elaboration *must* conclude with a list of the most significant and representative keywords and key terms/phrases directly identifiable from the transcribed document text.
                * This list is vital for the subsequent comparison process. Choose keywords that truly encapsulate the core topics and unique identifiers of the document's transcribed content.
                * Format this list *exactly* as follows, on a new line after the main re-elaboration:
                    `KEYWORDS: keyword1, keyword2, keyword3, term phrase 4, ...`

            ***IMPORTANT: DO NOT include this character in your response: ` ***

            PDF DOCUMENT TRANSCRIPTION (from .pdf):
            {text}

            RE-ELABORATION:""",

            '.xlsx': """You are an Expert Data Analyst and Document Specialist. Your task is to create a faithful and structured textual re-elaboration of the following text, which is an exact Python transcription from an Excel spreadsheet (`.xlsx`). The goal is to make this re-elaboration as similar as possible to the original transcribed data and structure, while presenting it with slightly enhanced text flow and clarity for subsequent AI analysis. **Avoid significant summarization, data aggregation (unless explicitly stated in the transcription), or interpretation beyond what the transcribed text directly supports.**

            This re-elaboration is critically important: it will serve as the primary reference point (representing an older version of the spreadsheet) for a subsequent AI-driven comparison against a newer version to detect significant changes in data or structure.

            Given that the text originates from an `.xlsx` file, the transcription likely represents tabular data, cell contents, sheet names, or textual descriptions of data structures. Your re-elaboration should textually represent this data structure and content faithfully.

            Please structure your output as follows:

            1.  **Main Textual Re-elaboration (Structured, Flowing, but True to Original Transcription):**
                * If the transcription provides an explicit overview or purpose of the spreadsheet, you may state it briefly. Otherwise, focus on re-elaborating the transcribed data.
                * Present the data categories, variables, values, and their apparent relationships largely as they appear in the provided transcription. If the transcription implies a tabular structure (e.g., rows and columns of data), your textual re-elaboration should reflect this structure clearly, perhaps by describing sheets, columns, and then their corresponding data, or by using consistent formatting to delineate records or data points.
                * Ensure that specific numerical values, text entries, headers, and any notes from the transcription are accurately represented, maintaining their original context and detail.
                * You may *lightly* improve the flow of the textual representation of the data (e.g., by introducing phrases like "The column 'X' contains values such as Y, Z" or "For the record/item A, the details are B, C"). However, **prioritize fidelity to the transcribed data, its structure, and its values over narrative creation or data interpretation.** The output should be a clear textual mirror of the transcribed spreadsheet content.

            2.  **Keywords (Formatted List - Crucial for Comparison):**
                * The re-elaboration *must* conclude with a list of the most significant and representative keywords, key terms, main data categories, and important column/row headers (if present and clearly identifiable as such in the transcription).
                * This list is vital for the subsequent comparison process. Choose terms that truly encapsulate the core data elements and structural identifiers of the spreadsheet's transcribed content.
                * Format this list *exactly* as follows, on a new line after the main re-elaboration:
                    `KEYWORDS: keyword1, data category A, header B, term phrase 4, ...`

            ***IMPORTANT: DO NOT include this character in your response: ` ***

            SPREADSHEET CONTENT TRANSCRIPTION (from .xlsx):
            {text}

            RE-ELABORATION:""",

            '.pptx': """You are an Expert Presentation Analyst. Your task is to create a faithful and structured textual re-elaboration of the following text, which is an exact Python transcription from a PowerPoint presentation (`.pptx`). The goal is to make this re-elaboration as similar as possible to the original transcribed slide content in terms of information, phrasing, and structure, while ensuring a slightly enhanced text flow and clarity for subsequent AI analysis. **Avoid significant summarization, content alteration, or reinterpretation of the slide messages.**

            This re-elaboration is critically important: it will serve as the primary reference point (representing an older version of the presentation) for a subsequent AI-driven comparison against a newer version to detect significant content changes.

            Given that the text originates from a `.pptx` file, the transcription likely represents slide titles, bullet points, main text blocks, and possibly speaker notes. Your re-elaboration should reflect this slide-based structure and content.

            Please structure your output as follows:

            1.  **Main Textual Re-elaboration (Flowing by Slide/Section, but True to Original Transcription):**
                * If the transcription begins with an overall presentation objective, you may briefly state it. Otherwise, proceed to re-elaborate the slide-by-slide content.
                * Present the content as a sequence of information blocks, likely corresponding to individual slides or logical sections of the presentation, as suggested by the transcription. Clearly delineate these blocks if possible (e.g., by stating "Slide Title: X" or "On the topic of Y:").
                * For each block/slide, reproduce the key messages, bullet points, data points, and any textual descriptions of visuals (if present in the transcription) with high fidelity to the original phrasing and detail.
                * You may *lightly* improve the flow between points within a perceived slide or between closely related "slides" if the transcription is very fragmented. For example, you can use simple connecting phrases. However, **prioritize fidelity to the transcribed content of each conceptual slide/section, its order, and its level of detail over extensive narrative creation or rephrasing.** The output should clearly reflect the information as it might have appeared on the slides.

            2.  **Keywords (Formatted List - Crucial for Comparison):**
                * The re-elaboration *must* conclude with a list of the most significant and representative keywords, key terms, and main topics covered in the presentation, as directly identifiable from the transcription.
                * This list is vital for the subsequent comparison process. Choose keywords that truly encapsulate the core topics and unique identifiers of the presentation's transcribed content.
                * Format this list *exactly* as follows, on a new line after the main re-elaboration:
                    `KEYWORDS: keyword1, keyword2, main topic A, term phrase 4, ...`


            ***IMPORTANT: DO NOT include this character in your response: ` ***

            PRESENTATION TRANSCRIPTION (from .pptx):
            {text}

            RE-ELABORATION:"""
        }

        # Seleziona il prompt appropriato o usa quello di default
        prompt_template = prompts.get(file_type.lower())

        if not prompt_template:
            print(f"Unsupported file format: {file_type}")
            return None

        prompt = prompt_template.format(text=text)
        
        response = await asyncio.to_thread(
            model.generate_content,
            prompt
        )
        
        return response.text
    except Exception as e:
        print(f"Error during summary generation: {e}")
        return None

async def summarize_document(file_path: str) -> str:
    """Carica e riassume un documento."""
    if not os.path.exists(file_path):
        return f"Error: File {file_path} does not exist."

    print(f"Loading document: {file_path}...")
    text = await get_document_text(file_path)
    
    if not text:
        return "Unable to extract text from the document."

    print("Generating summary...")
    _, file_extension = os.path.splitext(file_path)
    summary = await summarize_text(text, file_extension)
    
    if not summary:
        return "Error during summary generation."

    return summary

async def main():
    import sys
    
    if len(sys.argv) != 2:
        print("Usage: python agent.py <file_path>")
        print("Supported formats: .docx, .pdf, .xlsx, .pptx")
        return

    file_path = sys.argv[1]
    summary = await summarize_document(file_path)
    embedding = model.encode(file_content).tolist()
    
    print("\nFINAL SUMMARY:")
    print(summary)
    print("-" * 30)

if __name__ == "__main__":
    asyncio.run(main())