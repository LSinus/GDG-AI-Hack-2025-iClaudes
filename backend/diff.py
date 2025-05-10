import os
import asyncio
import google.generativeai as genai
import PyPDF2
import docx2txt
from unstructured.partition.xlsx import partition_xlsx
from pptx import Presentation

# Configurazione Gemini
GOOGLE_API_KEY = "AIzaSyASlsh1jRxd2iZnM3E3FM1TlDZ4yarUUMU"
genai.configure(api_key=GOOGLE_API_KEY)

# Configurazione del modello
model = genai.GenerativeModel('gemini-2.0-flash')

async def extract_text_from_docx(file_path):
    """Estrae il testo da un file DOCX."""
    try:
        text = docx2txt.process(file_path)
        return text
    except Exception as e:
        print(f"Errore nell'estrazione del testo da DOCX: {e}")
        return None

async def extract_text_from_pdf(file_path):
    """Estrae il testo da un file PDF."""
    try:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    except Exception as e:
        print(f"Errore nell'estrazione del testo da PDF: {e}")
        return None

async def extract_text_from_xlsx(file_path):
    """Estrae il testo da un file XLSX."""
    try:
        elements = partition_xlsx(file_path)
        text = "\n".join([str(element) for element in elements])
        return text
    except Exception as e:
        print(f"Errore nell'estrazione del testo da XLSX: {e}")
        return None

async def extract_text_from_pptx(file_path):
    """Estrae il testo da un file PowerPoint."""
    try:
        prs = Presentation(file_path)
        text = []
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text.append("\n".join(slide_text))
        
        return "\n\n".join(text)
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return None

async def get_document_text(file_path):
    """Estrae il testo dal documento in base al suo tipo."""
    _, file_extension = os.path.splitext(file_path)
    
    if file_extension.lower() == '.docx':
        return await extract_text_from_docx(file_path)
    elif file_extension.lower() == '.pdf':
        return await extract_text_from_pdf(file_path)
    elif file_extension.lower() in ['.xlsx', '.xls']:
        return await extract_text_from_xlsx(file_path)
    elif file_extension.lower() == '.pptx':
        return await extract_text_from_pptx(file_path)
    else:
        print(f"Unsupported file format: {file_extension}")
        return None

async def compare_versions(current_text, previous_summary):
    """Confronta le due versioni del testo usando Gemini."""
    try:
        prompt = """You are an Expert Comparative Document Content Analyst, with a distinct ability to discern significant evolutions between different versions of a text, even when one of them is represented in summary form. Your primary goal is to determine if content changes are substantial enough to mark a significant evolution.

**Your analysis must be deeply qualitative and semantic, focusing on the conceptual meaning and thematic shifts, not on superficial word-counting or exact phrase matching.**

Based on your comprehensive qualitative analysis of the two texts provided below, your **final output must SOLELY be a single number: 0 or 1.**

* **Output 0 if:** Your analysis concludes that the documents are still substantially similar. The changes are minor, incremental, or do not significantly alter the core message, scope, or main topics of the document.
* **Output 1 if:** Your analysis concludes that the documents are very different. There are significant new themes introduced, or substantial portions of previous content/themes have been removed or radically altered, leading to a noteworthy evolution of the document.

Here are the two texts you need to analyze to arrive at this numerical judgment:

--- FULL TEXT OF THE NEW VERSION ---
{current_text}
--- END OF FULL TEXT OF THE NEW VERSION ---

--- SUMMARY OF THE PREVIOUS VERSION WITH KEYWORDS ---
{previous_summary}
--- END OF SUMMARY OF THE PREVIOUS VERSION WITH KEYWORDS ---

**GUIDING STEPS FOR YOUR INTERNAL QUALITATIVE ANALYSIS (to arrive at the 0 or 1 decision):**

1.  **Comprehensive Understanding:**
    * Thoroughly understand the core content, central themes, key arguments, and main conclusions of the `FULL TEXT OF THE NEW VERSION`.
    * Thoroughly understand the inferred content of the previous version by analyzing the `SUMMARY OF THE PREVIOUS VERSION WITH KEYWORDS`. Focus on its core concepts and arguments, using the provided keywords as direct pointers to the actual content of that prior version.

2.  **Conceptual and Semantic Comparative Analysis (The foundation for your 0 or 1 decision):**
    * **Identify Emergence of New Core Concepts:** Does the `FULL TEXT OF THE NEW VERSION` introduce major themes, pivotal arguments, critical data, or new conclusions that were not present or substantively hinted at in the previous version (as inferred from its summary and keywords)? The introduction of *significant, new conceptual pillars* would lean towards a '1'.
    * **Identify Substantial Modification or Omission of Previous Core Concepts:** Are concepts that were clearly central to the previous version (inferred from its summary and keywords) now absent, significantly diminished in importance, or fundamentally re-framed in the `FULL TEXT OF THE NEW VERSION`? The removal or radical alteration of *previously foundational concepts* would lean towards a '1'.
    * **Strategic Keyword Analysis:**
        * Investigate the presence, absence, or altered context of the `Keywords from the previous version` within the `FULL TEXT OF THE NEW VERSION`.
        * Is the reappearance of a keyword superficial, or is the associated core concept meaningfully expanded, reduced, or re-contextualized?
        * Does the absence of previously central keywords strongly indicate a deliberate shift in thematic focus?
    * **Assess Overall Thematic and Structural Shifts:** Beyond individual concepts, has the document's overarching narrative, primary focus, argumentative structure, or intended purpose undergone a *major transformation* between the version inferred from the previous summary and the new full version? Such a transformation would lean towards a '1'.

3.  **Determining "Significance" for your 0 or 1 Output:**
    * Your internal judgment of "significant" should be based on whether the changes collectively result in a document that conveys a substantially new or different core message, addresses fundamentally new areas, or omits previously crucial elements to an extent that alters its overall identity or purpose.
    * **Minor edits, rephrasing, small additions/deletions that don't alter the core conceptual framework should lead you to output 0.**
    * **Introduction of entirely new major sections, removal of core pillars of the previous version, a radical shift in conclusions, or a deep restructuring that changes the fundamental message should lead you to output 1.**

**Reminder: Your final delivered output is ONLY the number 0 or 1, representing your ultimate qualitative and semantic judgment based on the thorough analysis outlined above.** Avoid any other text in the output."""

        prompt = prompt.format(current_text=current_text, previous_summary=previous_summary)
        
        response = await asyncio.to_thread(
            model.generate_content,
            prompt
        )
        
        # Estrai solo il numero dalla risposta
        result = response.text.strip()
        if result in ['0', '1']:
            return int(result)
        else:
            print(f"Unexpected response format: {result}")
            return None
    except Exception as e:
        print(f"Error during comparison: {e}")
        return None

async def compare_document_versions(file_path: str, previous_summary: str) -> int:
    """Confronta la versione corrente del documento con una versione precedente."""
    if not os.path.exists(file_path):
        print(f"Error: File {file_path} does not exist.")
        return None

    print(f"Loading current document: {file_path}...")
    current_text = await get_document_text(file_path)
    
    if not current_text:
        print("Unable to extract text from the current document.")
        return None

    print("Comparing versions...")
    result = await compare_versions(current_text, previous_summary)
    
    if result is None:
        print("Error during comparison.")
        return None

    return result

async def main():
    import sys
    
    if len(sys.argv) != 3:
        print("Usage: python diff.py <file_path> <previous_summary>")
        print("Supported formats: .docx, .pdf, .xlsx, .pptx")
        return

    file_path = sys.argv[1]
    previous_summary = sys.argv[2]
    
    result = await compare_document_versions(file_path, previous_summary)
    
    if result is not None:
        print("\nCOMPARISON RESULT:", result)
        print("0 = No significant changes")
        print("1 = Significant changes detected")
    print("-" * 30)

if __name__ == "__main__":
    asyncio.run(main())
